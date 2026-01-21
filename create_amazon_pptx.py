"""
Amazon TV XPath Guide PPT 생성 스크립트
BestBuy 버전처럼 수집 데이터 예시 포함
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def set_cell_style(cell, text, font_size=10, bold=False, bg_color=None, align_center=True):
    """셀 스타일 설정"""
    cell.text = text
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.name = '맑은 고딕'
    if align_center:
        paragraph.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

def add_table_slide(prs, title, headers, data, col_widths=None):
    """테이블이 있는 슬라이드 추가"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = '맑은 고딕'
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Table
    rows = len(data) + 1
    cols = len(headers)

    table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(0.8), Inches(9.4), Inches(0.5)).table

    # Set column widths
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)

    # Header row
    header_color = RGBColor(0, 51, 102)
    for i, header in enumerate(headers):
        set_cell_style(table.cell(0, i), header, font_size=11, bold=True, bg_color=header_color)
        table.cell(0, i).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    for row_idx, row_data in enumerate(data, 1):
        bg = RGBColor(240, 248, 255) if row_idx % 2 == 0 else None
        for col_idx, cell_data in enumerate(row_data):
            align = col_idx < 2  # 첫 두 컬럼만 가운데 정렬
            set_cell_style(table.cell(row_idx, col_idx), str(cell_data), font_size=9, bg_color=bg, align_center=align)

    return slide

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== Slide 1: Title ==========
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Amazon TV XPath 가이드"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = '맑은 고딕'
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "신규 입사자를 위한 XPath 수정 가이드"
    p.font.size = Pt(20)
    p.font.name = '맑은 고딕'
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Samsung DX Retail Crawler Team"
    p.font.size = Pt(14)
    p.font.name = '맑은 고딕'
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

    # ========== Slide 2: 개요 ==========
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📋 개요"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = '맑은 고딕'
    p.font.color.rgb = RGBColor(0, 51, 102)

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    contents = [
        ("📁 설정 테이블", "amazon_tv_config (category='xpath')"),
        ("📄 대상 파일", "amazon_tv_dt1.py, amazon_tv_dt2.py"),
        ("🔧 XPath 그룹핑", "config_key에서 _숫자 제거하여 그룹화\n    예: price_1, price_2 → price 리스트로 묶임"),
        ("⚡ 우선순위", "priority 값이 낮을수록 먼저 시도 (1이 최우선)"),
        ("🔄 Fallback", "DB에 없으면 코드 내 기본 XPath 사용"),
    ]

    for i, (title, desc) in enumerate(contents):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"{title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = '맑은 고딕'
        p.space_after = Pt(5)

        p2 = tf.add_paragraph()
        p2.text = f"    {desc}"
        p2.font.size = Pt(14)
        p2.font.name = '맑은 고딕'
        p2.space_after = Pt(20)

    # ========== Slide 3: 메인/BSR 페이지 ==========
    headers = ["config_key", "설명", "수집 데이터 예시", "XPath 예시"]
    data = [
        ["product_title", "상품명", "Samsung 65\" Class OLED S90D", "//span[@id='productTitle']"],
        ["product_price", "가격", "$1,597.99", "//span[@class='a-price-whole']"],
        ["product_image", "상품 이미지 URL", "https://m.media-amazon.com/images/I/...", "//img[@id='landingImage']/@src"],
        ["product_asin", "ASIN (상품코드)", "B0CTS8P319", "//div[@data-asin]/@data-asin"],
        ["product_rating", "평점", "4.5", "//span[@class='a-icon-alt']"],
        ["review_count", "리뷰 수", "1,234", "//span[@id='acrCustomerReviewText']"],
    ]
    add_table_slide(prs, "📦 메인/BSR 페이지 XPath", headers, data, [1.5, 1.5, 3, 3.4])

    # ========== Slide 4: 상세 페이지 - 기본 정보 ==========
    headers = ["config_key", "설명", "수집 데이터 예시", "XPath 예시"]
    data = [
        ["screen_size", "화면 크기", "65 inches", "//tr[contains(.,'Screen Size')]//td"],
        ["model_name", "모델명", "QN65S90DAFXZA", "//tr[contains(.,'Model Name')]//td"],
        ["display_tech", "디스플레이 기술", "OLED", "//tr[contains(.,'Display Technology')]//td"],
        ["resolution", "해상도", "4K", "//tr[contains(.,'Resolution')]//td"],
        ["brand", "브랜드", "SAMSUNG", "//tr[contains(.,'Brand')]//td"],
        ["refresh_rate", "주사율", "120 Hz", "//tr[contains(.,'Refresh Rate')]//td"],
    ]
    add_table_slide(prs, "📱 상세 페이지 - 기본 정보", headers, data, [1.5, 1.5, 3, 3.4])

    # ========== Slide 5: 상세 페이지 - 가격 정보 ==========
    headers = ["config_key", "설명", "수집 데이터 예시", "XPath 예시"]
    data = [
        ["final_price", "최종 판매가", "$1,597.99", "//span[@class='a-price']//span[@class='a-offscreen']"],
        ["original_price", "원래 가격 (정가)", "$2,499.99", "//span[@class='a-price a-text-price']//span"],
        ["price_savings", "할인 금액", "Save $902.00", "//span[contains(@class,'savingsPercentage')]"],
        ["price_status", "가격 상태", "In Stock", "//div[@id='availability']//span"],
        ["sold_by", "판매자", "Amazon.com", "//div[@id='merchant-info']"],
        ["ships_from", "배송 출발지", "Amazon", "//div[@id='fulfilledBy']"],
    ]
    add_table_slide(prs, "💰 상세 페이지 - 가격 정보", headers, data, [1.5, 1.5, 3, 3.4])

    # ========== Slide 6: 상세 페이지 - 리뷰 정보 ==========
    headers = ["config_key", "설명", "수집 데이터 예시", "XPath 예시"]
    data = [
        ["star_rating", "별점 (5점 만점)", "4.5 out of 5 stars", "//span[@id='acrPopover']/@title"],
        ["count_of_reviews", "전체 리뷰 수", "1,234 ratings", "//span[@id='acrCustomerReviewText']"],
        ["review_link", "리뷰 페이지 링크", "/product-reviews/B0CTS8P319", "//a[@data-hook='see-all-reviews-link']"],
        ["5_star_pct", "5점 비율", "70%", "//tr[@class='a-histogram-row'][1]//td[3]"],
        ["4_star_pct", "4점 비율", "15%", "//tr[@class='a-histogram-row'][2]//td[3]"],
        ["1_star_pct", "1점 비율", "5%", "//tr[@class='a-histogram-row'][5]//td[3]"],
    ]
    add_table_slide(prs, "⭐ 상세 페이지 - 리뷰 정보", headers, data, [1.5, 1.5, 3, 3.4])

    # ========== Slide 7: 자주 변경하는 XPath ==========
    headers = ["config_key", "변경 빈도", "변경 원인", "확인 방법"]
    data = [
        ["final_price", "높음", "Amazon 가격 표시 구조 변경", "상품 상세 페이지에서 F12 → 가격 요소 확인"],
        ["star_rating", "중간", "리뷰 섹션 HTML 변경", "별점 아이콘 주변 요소 검사"],
        ["sold_by", "중간", "판매자 정보 영역 변경", "Add to Cart 버튼 하단 확인"],
        ["screen_size", "낮음", "상품 스펙 테이블 구조 안정적", "Product Details 섹션 확인"],
        ["product_image", "낮음", "이미지 ID 안정적", "메인 이미지 요소 확인"],
    ]
    add_table_slide(prs, "🔄 자주 변경되는 XPath 목록", headers, data, [2, 1.5, 3, 3])

    # ========== Slide 8: XPath 수정 예시 ==========
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🛠️ XPath 수정 방법"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = '맑은 고딕'
    p.font.color.rgb = RGBColor(0, 51, 102)

    content_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.9), Inches(9.4), Inches(6))
    tf = content_box.text_frame
    tf.word_wrap = True

    steps = [
        ("1️⃣ 문제 확인", "크롤링 결과에서 특정 필드가 NULL이거나 잘못된 값인지 확인"),
        ("2️⃣ 브라우저에서 확인", "Amazon 페이지 → F12 개발자도구 → 해당 요소 찾기\n    Elements 탭에서 Ctrl+F로 XPath 테스트"),
        ("3️⃣ DB에서 XPath 수정", """SQL 예시:
    UPDATE amazon_tv_config
    SET config_value = '//새로운/xpath'
    WHERE config_key = 'final_price_1'
    AND category = 'xpath';"""),
        ("4️⃣ 새 XPath 추가 (Fallback)", """SQL 예시:
    INSERT INTO amazon_tv_config
    (category, config_key, config_value, priority, is_active)
    VALUES ('xpath', 'final_price_2', '//대체/xpath', 2, TRUE);"""),
        ("5️⃣ 테스트", "크롤러 실행하여 데이터 수집 확인"),
    ]

    for i, (title, desc) in enumerate(steps):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = '맑은 고딕'
        p.space_after = Pt(3)

        p2 = tf.add_paragraph()
        p2.text = f"    {desc}"
        p2.font.size = Pt(11)
        p2.font.name = 'Consolas' if 'SQL' in desc or 'xpath' in desc else '맑은 고딕'
        p2.space_after = Pt(12)

    # ========== Slide 9: 주의사항 ==========
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "⚠️ 주의사항"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = '맑은 고딕'
    p.font.color.rgb = RGBColor(0, 51, 102)

    content_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.9), Inches(9.4), Inches(6))
    tf = content_box.text_frame
    tf.word_wrap = True

    warnings = [
        ("🔴 is_active = TRUE 확인", "비활성화된 XPath는 로드되지 않음"),
        ("🔴 priority 순서 주의", "숫자가 낮을수록 먼저 시도됨 (1 > 2 > 3)"),
        ("🔴 특수문자 이스케이프", "XPath 내 따옴표 주의: ' → '' (SQL에서)"),
        ("🔴 테스트 환경에서 먼저 확인", "운영 DB 직접 수정 전 로컬 테스트 필수"),
        ("🔴 기존 XPath 삭제 금지", "is_active = FALSE로 변경 후 새 XPath 추가 권장"),
        ("🔴 config_key 네이밍 규칙", "기존 패턴 유지: {필드명}_{우선순위번호}\n    예: final_price_1, final_price_2"),
    ]

    for i, (title, desc) in enumerate(warnings):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = '맑은 고딕'
        p.space_after = Pt(3)

        p2 = tf.add_paragraph()
        p2.text = f"    → {desc}"
        p2.font.size = Pt(13)
        p2.font.name = '맑은 고딕'
        p2.space_after = Pt(15)

    # Save
    output_path = r'c:\Users\gomguard\Documents\퀵오일\삼성전자\samsung_dx_retail_com\samsung_dx_retail_com\amazon_tv_xpath_guide.pptx'
    prs.save(output_path)
    print(f"[OK] PPT created: {output_path}")
    return output_path

if __name__ == '__main__':
    create_pptx()
