from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=''):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, content_lines):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    return slide

def add_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    num_rows = len(rows) + 1
    num_cols = len(headers)

    left = Inches(0.3)
    top = Inches(1.1)
    width = Inches(12.7)
    height = Inches(0.4 * num_rows)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    col_widths = [Inches(12.7 / num_cols)] * num_cols
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    return slide

# 1. 표지
add_title_slide(prs, 'BestBuy TV 크롤러', 'XPath 설정 가이드 (신규입사자용)')

# 2. 개요
add_content_slide(prs, '1. 개요', [
    '■ 설정 테이블: bby_tv_config',
    '',
    '■ XPath 수정 SQL:',
    "   UPDATE bby_tv_config",
    "   SET config_value = '새로운_xpath'",
    "   WHERE category = 'xpath'",
    "     AND config_key = '컬럼명'",
    "     AND file_name = '파일명'",
    "     AND priority = 1;",
    '',
    '■ priority: 1순위 실패 시 2순위 시도 (fallback 구조)'
])

# 3. 리스트 페이지
add_table_slide(prs, '2. 리스트 페이지 (bby_tv_main1, bby_tv_bsr1)',
    ['config_key', '설명', '수집 데이터 예시'],
    [
        ['product_container', '상품 카드 전체 영역', '<li class="product-list-item">'],
        ['product_title', '상품명', 'Samsung 65" Class OLED S90D'],
        ['product_url', '상품 상세 URL', '/site/samsung-65.../6576038.p'],
        ['product_link', '상품 링크 요소', '클릭 가능한 <a> 태그'],
        ['offer', '추가 할인 정보', '+4 Offers'],
        ['pickup_availability', '매장 픽업 가능 여부', 'Pick up today'],
        ['shipping_availability', '배송 정보', 'Get it by Jan 25'],
        ['delivery_availability', '배달 가능 여부', 'Delivery available'],
        ['sponsored', '광고 상품 여부', '광고 배지 유무'],
    ])

# 4. 프로모션 페이지
add_table_slide(prs, '3. 프로모션 페이지 (bby_tv_pmt1)',
    ['config_key', '설명', '수집 데이터 예시'],
    [
        ['all_sections', '프로모션 섹션 전체', '<section> 영역들'],
        ['carousel_list', '상품 캐러셀 리스트', '가로 스크롤 상품 목록'],
        ['product_item', '개별 상품 아이템', '캐러셀 내 각 상품'],
        ['product_name', '상품명', 'LG 55" Class OLED B4'],
        ['product_url', '상품 URL', '상세 페이지 링크'],
        ['offer', '할인 정보', 'Save $300'],
        ['h2_headline', '섹션 제목 (H2)', 'Top TV Deals'],
        ['p_subtitle', '섹션 부제목', 'Limited time offers'],
        ['section_title', '섹션 타이틀', '각 프로모션 영역 제목'],
    ])

# 5. 트렌딩 딜
add_table_slide(prs, '4. 트렌딩 딜 (bby_tv_trend_crawl)',
    ['config_key', '설명', '수집 데이터 예시'],
    [
        ['trending_section', '트렌딩 딜 섹션 영역', 'Trending deals 컨테이너'],
        ['tvs_button', 'TVs & Projectors 탭 버튼', '카테고리 선택 버튼'],
        ['product_items', '상품 목록 컨테이너', '트렌딩 상품들'],
        ['product_name', '상품명', 'Sony 75" Class BRAVIA'],
        ['product_url', '상품 URL', '상세 페이지 링크'],
        ['rank', '트렌딩 순위', '#1, #2 등'],
    ])

# 6. 상세 페이지
add_table_slide(prs, '5. 상세 페이지 (bby_tv_dt1)',
    ['config_key', '설명', '수집 데이터 예시'],
    [
        ['product_title', '상품명 (상세)', '전체 상품명'],
        ['final_price', '최종 판매가', '$1,299.99'],
        ['model_number', '모델 번호', 'QN65S90D'],
        ['overall_rating', '평점', '4.7 out of 5'],
        ['review_count', '리뷰 수', '(1,234 reviews)'],
        ['specs_button', '스펙 펼치기 버튼', '상세 스펙 보기'],
        ['see_all_reviews_btn', '리뷰 더보기 버튼', '전체 리뷰 페이지 이동'],
        ['next_page_btn', '리뷰 다음 페이지', '리뷰 페이지네이션'],
        ['top_mentions', '리뷰 키워드', 'Picture quality, Easy setup'],
    ])

# 7. 자주 수정하는 XPath
add_table_slide(prs, '6. 자주 수정하는 XPath',
    ['순위', 'config_key', '파일', '변경 빈도', '사유'],
    [
        ['1', 'product_title', '전체', '높음', '클래스명 변경'],
        ['2', 'product_url', '전체', '높음', 'data-testid 변경'],
        ['3', 'tvs_button', 'trend_crawl', '높음', '홈페이지 구조 변경'],
        ['4', 'top_mentions', 'dt1', '중간', '리뷰 섹션 개편'],
        ['5', 'offer', 'main1, pmt1', '중간', '할인 배지 디자인 변경'],
    ])

# 8. 수정 예시
add_content_slide(prs, '7. XPath 수정 예시', [
    '■ 상황: 상품명 클래스가 product-title → item-name으로 변경됨',
    '',
    '-- 1순위 XPath 수정',
    'UPDATE bby_tv_config',
    "SET config_value = './/h2[contains(@class, \"item-name\")]'",
    "WHERE category = 'xpath'",
    "  AND config_key = 'product_title'",
    "  AND file_name = 'bby_tv_main1'",
    '  AND priority = 1;',
])

# 9. 주의사항
add_content_slide(prs, '8. 주의사항', [
    '1. 테스트 필수',
    '   - 수정 후 해당 크롤러 단독 실행하여 확인',
    '',
    '2. 백업',
    '   - 수정 전 기존 값 메모',
    '   - 또는 is_active = FALSE로 비활성화 후 새로 추가',
    '',
    '3. priority 순서',
    '   - 낮은 숫자가 먼저 시도됨 (1 → 2 → 3)',
    '',
    '4. file_name 정확히',
    '   - 오타 시 적용 안 됨',
])

# 저장
output_path = 'BestBuy_TV_XPath_Guide.pptx'
prs.save(output_path)
print(f'[OK] 파일 생성 완료: {output_path}')
