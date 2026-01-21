from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.util import Pt
import psycopg2
from config import DB_CONFIG

from pptx.dml.color import RGBColor

def RgbColor(r, g, b):
    """RGB 색상 반환"""
    return RGBColor(r, g, b)

# 색상 정의 (삼성 블루 계열)
SAMSUNG_BLUE = RgbColor(0x12, 0x28, 0x65)  # 삼성 다크블루
ACCENT_BLUE = RgbColor(0x00, 0x71, 0xC5)   # 밝은 블루
LIGHT_BLUE = RgbColor(0xE8, 0xF4, 0xFC)    # 연한 블루 배경
DARK_GRAY = RgbColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RgbColor(0xF5, 0xF5, 0xF5)
WHITE = RgbColor(0xFF, 0xFF, 0xFF)
RED = RgbColor(0xE5, 0x3E, 0x3E)
GREEN = RgbColor(0x2E, 0xA0, 0x4D)
ORANGE = RgbColor(0xF5, 0xA6, 0x23)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_cell_fill(cell, color):
    """셀 배경색 설정 - color는 (r, g, b) 튜플"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # 기존 fill 제거
    for child in tcPr:
        if 'solidFill' in child.tag:
            tcPr.remove(child)
    solidFill = OxmlElement('a:solidFill')
    srgbClr = OxmlElement('a:srgbClr')
    srgbClr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def add_header_bar(slide, title_text):
    """상단 헤더 바 추가"""
    # 헤더 배경
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SAMSUNG_BLUE
    shape.line.fill.background()

    # 헤더 텍스트
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_footer(slide, page_num, total_pages):
    """하단 푸터 추가"""
    # 푸터 라인
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(13.333), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    # 페이지 번호
    txBox = slide.shapes.add_textbox(Inches(12), Inches(7.25), Inches(1), Inches(0.25))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num} / {total_pages}"
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.RIGHT

    # 문서 제목
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.25), Inches(5), Inches(0.25))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "BestBuy TV Crawler - XPath 설정 가이드"
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY


def add_title_slide(prs, title, subtitle):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경 그라데이션 효과 (상단 바)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SAMSUNG_BLUE
    shape.line.fill.background()

    # 하단 악센트 바
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(3), Inches(13.333), Inches(0.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    # 메인 타이틀
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(0xAA, 0xCC, 0xEE)
    p.alignment = PP_ALIGN.CENTER

    # 하단 정보
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Samsung DS Retail Monitoring Team"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "2025년 1월"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_toc_slide(prs, page_num, total_pages):
    """목차 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "목차")
    add_footer(slide, page_num, total_pages)

    toc_items = [
        ("01", "개요 및 설정 방법", "설정 테이블 구조, 수정 SQL"),
        ("02", "크롤러 구조도", "전체 크롤러 실행 흐름"),
        ("03", "리스트 페이지 XPath", "bby_tv_main1, bby_tv_bsr1"),
        ("04", "프로모션 페이지 XPath", "bby_tv_pmt1"),
        ("05", "트렌딩 딜 XPath", "bby_tv_trend_crawl"),
        ("06", "상세 페이지 XPath", "bby_tv_dt1"),
        ("07", "자주 변경되는 XPath", "변경 빈도 및 사유"),
        ("08", "XPath 수정 예시", "실제 SQL 예시"),
        ("09", "Quick Reference", "한눈에 보는 요약"),
    ]

    start_y = 1.2
    for i, (num, title, desc) in enumerate(toc_items):
        y = start_y + (i * 0.65)

        # 번호 박스
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(0.6), Inches(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE
        shape.line.fill.background()

        # 번호 텍스트
        tf = shape.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 제목
        txBox = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(5), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY

        # 설명
        txBox = slide.shapes.add_textbox(Inches(1.3), Inches(y + 0.3), Inches(5), Inches(0.25))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(0x66, 0x66, 0x66)

    return slide


def add_overview_slide(prs, page_num, total_pages):
    """개요 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "01. 개요 및 설정 방법")
    add_footer(slide, page_num, total_pages)

    # 왼쪽: 설정 테이블 구조
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(6), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "📋 설정 테이블 구조"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SAMSUNG_BLUE

    # 테이블 구조 설명
    table_info = [
        ("테이블명", "bby_tv_config"),
        ("category", "xpath, url, timing, constant, table, retry"),
        ("config_key", "설정 키 (예: product_title)"),
        ("config_value", "실제 값 (XPath, URL 등)"),
        ("file_name", "적용 파일 (예: bby_tv_main1)"),
        ("priority", "우선순위 (1이 최우선)"),
        ("is_active", "활성화 여부 (TRUE/FALSE)"),
    ]

    table = slide.shapes.add_table(len(table_info), 2, Inches(0.5), Inches(1.6), Inches(5.5), Inches(2.8)).table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(3.7)

    for i, (col1, col2) in enumerate(table_info):
        cell0 = table.cell(i, 0)
        cell0.text = col1
        cell0.text_frame.paragraphs[0].font.size = Pt(12)
        cell0.text_frame.paragraphs[0].font.bold = True
        set_cell_fill(cell0, (0xE8, 0xF4, 0xFC))

        cell1 = table.cell(i, 1)
        cell1.text = col2
        cell1.text_frame.paragraphs[0].font.size = Pt(11)

    # 오른쪽: 수정 SQL
    txBox = slide.shapes.add_textbox(Inches(6.8), Inches(1.1), Inches(6), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "🔧 XPath 수정 SQL"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SAMSUNG_BLUE

    # SQL 코드 박스
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(6), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0x2D, 0x2D, 0x2D)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].space_before = Pt(10)

    sql_lines = [
        "UPDATE bby_tv_config",
        "SET config_value = '새로운_xpath'",
        "WHERE category = 'xpath'",
        "  AND config_key = 'product_title'",
        "  AND file_name = 'bby_tv_main1'",
        "  AND priority = 1;",
    ]

    for i, line in enumerate(sql_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "  " + line
        p.font.size = Pt(13)
        p.font.name = "Consolas"
        p.font.color.rgb = RgbColor(0x50, 0xFA, 0x7B) if line.startswith("UPDATE") or line.startswith("SET") else WHITE

    # 하단: Priority 설명
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "⚡ Priority (우선순위) 작동 방식"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SAMSUNG_BLUE

    # Priority 다이어그램
    priorities = [
        ("Priority 1", "1순위 시도", ACCENT_BLUE),
        ("→", "", DARK_GRAY),
        ("실패 시", "", RED),
        ("→", "", DARK_GRAY),
        ("Priority 2", "2순위 시도", RgbColor(0x5B, 0xA0, 0xD0)),
        ("→", "", DARK_GRAY),
        ("실패 시", "", RED),
        ("→", "", DARK_GRAY),
        ("Priority 3", "3순위 시도", RgbColor(0x8B, 0xC0, 0xE0)),
    ]

    x = 0.5
    for label, desc, color in priorities:
        if label == "→":
            txBox = slide.shapes.add_textbox(Inches(x), Inches(5.2), Inches(0.4), Inches(0.4))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "→"
            p.font.size = Pt(20)
            p.font.color.rgb = color
            x += 0.5
        elif label == "실패 시":
            txBox = slide.shapes.add_textbox(Inches(x), Inches(5.25), Inches(0.8), Inches(0.3))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "실패 시"
            p.font.size = Pt(11)
            p.font.color.rgb = color
            x += 0.9
        else:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.15), Inches(1.4), Inches(0.6)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()

            tf = shape.text_frame
            tf.paragraphs[0].text = label
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = WHITE
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            x += 1.6

    return slide


def add_architecture_slide(prs, page_num, total_pages):
    """크롤러 구조도 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "02. 크롤러 구조도")
    add_footer(slide, page_num, total_pages)

    # 메인 크롤러 (bby_tv_crawl.py)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(1.2), Inches(3.333), Inches(0.7)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SAMSUNG_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = "bby_tv_crawl.py"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 화살표들
    crawlers = [
        ("bby_tv_main1.py", "메인 TV 목록", Inches(0.5), "bby_tv_main1"),
        ("bby_tv_bsr1.py", "베스트셀러", Inches(3.2), "bby_tv_bsr1"),
        ("bby_tv_pmt1.py", "프로모션", Inches(5.9), "bby_tv_pmt1"),
        ("bby_tv_trend_crawl.py", "트렌딩 딜", Inches(8.6), "bby_tv_trend_crawl"),
        ("bby_tv_dt1.py", "상세 페이지", Inches(11.3), "bby_tv_dt1"),
    ]

    for name, desc, x, table in crawlers:
        # 연결선
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x + Inches(0.9), Inches(1.9), Inches(0.03), Inches(0.5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_BLUE
        line.line.fill.background()

        # 크롤러 박스
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.4), Inches(1.8), Inches(1.1)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = name.replace(".py", "")
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = RgbColor(0xCC, 0xE5, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 화살표 (dt1으로)
    txBox = slide.shapes.add_textbox(Inches(4.5), Inches(3.6), Inches(6), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "URL 수집 ─────────────────────────────────────────→"
    p.font.size = Pt(12)
    p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER

    # DB 저장 테이블
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "📦 저장 테이블"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SAMSUNG_BLUE

    tables = [
        ("bby_tv_main1", "메인 목록"),
        ("bby_tv_bsr1", "베스트셀러"),
        ("bby_tv_pmt1", "프로모션"),
        ("bby_tv_Trend_crawl", "트렌딩"),
        ("bby_tv_crawl", "상세 데이터"),
        ("tv_retail_com", "통합 데이터"),
    ]

    x = 0.5
    for tbl_name, tbl_desc in tables:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.8), Inches(2), Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(0x4C, 0xAF, 0x50)
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = tbl_name
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = tbl_desc
        p.font.size = Pt(9)
        p.font.color.rgb = RgbColor(0xCC, 0xFF, 0xCC)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        x += 2.1

    # 설정 테이블 (우측)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10), Inches(5.8), Inches(2.8), Inches(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

    tf = shape.text_frame
    tf.paragraphs[0].text = "bby_tv_config"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "XPath, URL, Timing 설정"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 화살표
    txBox = slide.shapes.add_textbox(Inches(8), Inches(6.1), Inches(2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "설정 로드 ←"
    p.font.size = Pt(11)
    p.font.color.rgb = ORANGE

    return slide


def add_xpath_table_slide(prs, title, subtitle, file_name, headers, rows, page_num, total_pages):
    """XPath 테이블 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, title)
    add_footer(slide, page_num, total_pages)

    # 서브타이틀
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"📁 파일: {file_name}  |  {subtitle}"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(0x66, 0x66, 0x66)

    # 테이블
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.3), Inches(1.6), Inches(12.7), Inches(5.2)).table

    # 컬럼 너비 설정
    if num_cols == 3:
        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(3.5)
        table.columns[2].width = Inches(6.7)
    elif num_cols == 4:
        table.columns[0].width = Inches(2.2)
        table.columns[1].width = Inches(2.5)
        table.columns[2].width = Inches(1.3)
        table.columns[3].width = Inches(6.7)

    # 헤더
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(13)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        set_cell_fill(cell, (0x12, 0x28, 0x65))

    # 데이터
    for row_idx, row_data in enumerate(rows):
        bg_color = (0xF5, 0xF5, 0xF5) if row_idx % 2 == 0 else (0xFF, 0xFF, 0xFF)
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            if col_idx == num_cols - 1:  # XPath 컬럼
                cell.text_frame.paragraphs[0].font.name = "Consolas"
                cell.text_frame.paragraphs[0].font.size = Pt(10)
            set_cell_fill(cell, bg_color)

    return slide


def add_frequency_slide(prs, page_num, total_pages):
    """자주 변경되는 XPath 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "07. 자주 변경되는 XPath")
    add_footer(slide, page_num, total_pages)

    # 변경 빈도 높음
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(6), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "🔴 변경 빈도: 높음"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED

    high_freq = [
        ("product_title", "전체", "클래스명 변경이 잦음"),
        ("product_url", "전체", "data-testid 속성 변경"),
        ("tvs_button", "trend_crawl", "홈페이지 리뉴얼 시 구조 변경"),
    ]

    table = slide.shapes.add_table(4, 3, Inches(0.5), Inches(1.6), Inches(5.8), Inches(1.8)).table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(2.5)

    headers = ["config_key", "파일", "변경 사유"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        set_cell_fill(cell, (0xE5, 0x3E, 0x3E))
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE

    for row_idx, row_data in enumerate(high_freq):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(11)

    # 변경 빈도 중간
    txBox = slide.shapes.add_textbox(Inches(6.8), Inches(1.1), Inches(6), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "🟡 변경 빈도: 중간"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    mid_freq = [
        ("top_mentions", "dt1", "리뷰 UI 개편"),
        ("offer", "main1, pmt1", "할인 배지 디자인"),
        ("carousel_list", "pmt1", "프로모션 레이아웃"),
    ]

    table = slide.shapes.add_table(4, 3, Inches(6.8), Inches(1.6), Inches(5.8), Inches(1.8)).table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(2.5)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        set_cell_fill(cell, (0xF5, 0xA6, 0x23))
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE

    for row_idx, row_data in enumerate(mid_freq):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(11)

    # 변경 빈도 낮음
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "🟢 변경 빈도: 낮음 (거의 변경 없음)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN

    low_freq = [
        ("product_container", "전체", "상품 카드 구조 - 대규모 리뉴얼 외 안정적"),
        ("trending_section", "trend_crawl", "섹션 컨테이너 - 구조적으로 안정"),
        ("specs_button", "dt1", "스펙 버튼 - 상세페이지 기본 요소"),
        ("next_page_btn", "dt1", "페이지네이션 - 표준 UI 컴포넌트"),
    ]

    table = slide.shapes.add_table(5, 3, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.3)).table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(8.3)

    for i, h in enumerate(["config_key", "파일", "안정적인 이유"]):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        set_cell_fill(cell, (0x2E, 0xA0, 0x4D))
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE

    for row_idx, row_data in enumerate(low_freq):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(11)

    return slide


def add_example_slide(prs, page_num, total_pages):
    """수정 예시 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "08. XPath 수정 예시")
    add_footer(slide, page_num, total_pages)

    # 시나리오
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "📌 시나리오: 상품명 클래스가 'product-title' → 'item-name'으로 변경됨"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    # 방법 1
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(6), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "방법 1: 기존 XPath 직접 수정"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2), Inches(6), Inches(2.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0x2D, 0x2D, 0x2D)
    shape.line.fill.background()

    sql1 = """UPDATE bby_tv_config
SET config_value =
  './/h2[contains(@class, "item-name")]'
WHERE category = 'xpath'
  AND config_key = 'product_title'
  AND file_name = 'bby_tv_main1'
  AND priority = 1;"""

    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(sql1.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "  " + line
        p.font.size = Pt(12)
        p.font.name = "Consolas"
        p.font.color.rgb = WHITE

    # 방법 2
    txBox = slide.shapes.add_textbox(Inches(6.8), Inches(1.7), Inches(6), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "방법 2: 새 XPath 추가 (fallback 유지)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(6), Inches(2.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0x2D, 0x2D, 0x2D)
    shape.line.fill.background()

    sql2 = """-- 기존 priority 2로 변경
UPDATE bby_tv_config
SET priority = 2
WHERE config_key = 'product_title'
  AND file_name = 'bby_tv_main1'
  AND priority = 1;

-- 새 XPath를 priority 1로 추가
INSERT INTO bby_tv_config
  (category, config_key, config_value,
   file_name, priority, is_active)
VALUES ('xpath', 'product_title',
  './/h2[contains(@class, "item-name")]',
  'bby_tv_main1', 1, TRUE);"""

    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(sql2.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = " " + line
        p.font.size = Pt(9)
        p.font.name = "Consolas"
        if line.startswith("--"):
            p.font.color.rgb = RgbColor(0x6A, 0x99, 0x55)
        else:
            p.font.color.rgb = WHITE

    # 확인 방법
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "✅ 수정 후 확인"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GREEN

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0x2D, 0x2D, 0x2D)
    shape.line.fill.background()

    verify = """-- 설정 확인
SELECT config_key, priority, config_value FROM bby_tv_config
WHERE config_key = 'product_title' AND file_name = 'bby_tv_main1' ORDER BY priority;

-- 테스트 실행 (단일 크롤러)
python bby_tv_main1.py"""

    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(verify.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "  " + line
        p.font.size = Pt(11)
        p.font.name = "Consolas"
        if line.startswith("--"):
            p.font.color.rgb = RgbColor(0x6A, 0x99, 0x55)
        else:
            p.font.color.rgb = WHITE

    return slide


def add_quick_ref_slide(prs, page_num, total_pages):
    """Quick Reference 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "09. Quick Reference")
    add_footer(slide, page_num, total_pages)

    # 파일별 주요 XPath
    files_info = [
        ("bby_tv_main1", "메인 목록", ["product_title", "product_url", "offer"]),
        ("bby_tv_bsr1", "베스트셀러", ["product_title", "product_url", "offer"]),
        ("bby_tv_pmt1", "프로모션", ["product_name", "product_url", "carousel_list"]),
        ("bby_tv_trend_crawl", "트렌딩", ["tvs_button", "product_name", "rank"]),
        ("bby_tv_dt1", "상세", ["product_title", "final_price", "top_mentions"]),
    ]

    x = 0.3
    for file_name, desc, keys in files_info:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.2), Inches(2.4), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BLUE
        shape.line.color.rgb = ACCENT_BLUE

        tf = shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = file_name.replace("bby_tv_", "")
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = SAMSUNG_BLUE
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = "─" * 10
        p.font.size = Pt(8)
        p.font.color.rgb = ACCENT_BLUE
        p.alignment = PP_ALIGN.CENTER

        for key in keys:
            p = tf.add_paragraph()
            p.text = f"• {key}"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY

        x += 2.55

    # 주의사항 박스
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(4), Inches(12.7), Inches(2.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0xFF, 0xF3, 0xE0)
    shape.line.color.rgb = ORANGE

    tf = shape.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "⚠️ 주의사항"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    warnings = [
        "1. 수정 전 기존 값 백업 (메모 또는 is_active=FALSE 처리)",
        "2. 수정 후 반드시 단일 크롤러로 테스트",
        "3. file_name 오타 주의 (적용 안 됨)",
        "4. priority 숫자가 낮을수록 먼저 시도",
        "5. 문제 발생 시 is_active=FALSE로 비활성화 후 롤백",
    ]

    for w in warnings:
        p = tf.add_paragraph()
        p.text = w
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY

    return slide


# === 메인 실행 ===
print("[INFO] DB에서 XPath 정보 로드 중...")

# DB에서 실제 XPath 데이터 가져오기
conn = psycopg2.connect(**DB_CONFIG)
cursor = conn.cursor()

cursor.execute("""
    SELECT file_name, config_key, config_value, priority
    FROM bby_tv_config
    WHERE category = 'xpath' AND is_active = TRUE
    ORDER BY file_name, config_key, priority
""")
xpath_data = cursor.fetchall()
cursor.close()
conn.close()

# 파일별로 그룹화
from collections import defaultdict
by_file = defaultdict(list)
for row in xpath_data:
    by_file[row[0]].append((row[1], row[2], row[3]))

total_pages = 11

# 1. 표지
add_title_slide(prs, "BestBuy TV 크롤러", "XPath 설정 가이드 (신규입사자용)")

# 2. 목차
add_toc_slide(prs, 2, total_pages)

# 3. 개요
add_overview_slide(prs, 3, total_pages)

# 4. 크롤러 구조도
add_architecture_slide(prs, 4, total_pages)

# 5. 리스트 페이지 XPath (main1, bsr1)
main_rows = []
for key, val, pri in by_file.get('bby_tv_main1', []):
    desc = {
        'product_container': '상품 카드 전체 영역',
        'product_title': '상품명',
        'product_url': '상품 URL',
        'product_link': '상품 링크 요소',
        'offer': '추가 할인 정보',
        'pickup_availability': '매장 픽업',
        'shipping_availability': '배송 정보',
        'delivery_availability': '배달 가능',
        'sponsored': '광고 여부',
    }.get(key, key)
    # XPath 길이 제한
    val_display = val[:55] + "..." if len(val) > 55 else val
    main_rows.append((key, desc, f"p{pri}", val_display))

add_xpath_table_slide(prs, "03. 리스트 페이지 XPath",
    "메인 TV 목록 / 베스트셀러 목록 (동일 구조)",
    "bby_tv_main1.py, bby_tv_bsr1.py",
    ["config_key", "설명", "우선순위", "XPath"],
    main_rows, 5, total_pages)

# 6. 프로모션 페이지 XPath
pmt_rows = []
for key, val, pri in by_file.get('bby_tv_pmt1', []):
    desc = {
        'all_sections': '프로모션 섹션 전체',
        'carousel_list': '상품 캐러셀',
        'product_item': '개별 상품',
        'product_name': '상품명',
        'product_url': '상품 URL',
        'offer': '할인 정보',
        'h2_headline': '섹션 제목',
        'p_subtitle': '섹션 부제목',
        'section_title': '섹션 타이틀',
        'promo_type_hero': '히어로 배너',
        'section_p': '섹션 설명',
    }.get(key, key)
    val_display = val[:55] + "..." if len(val) > 55 else val
    pmt_rows.append((key, desc, f"p{pri}", val_display))

add_xpath_table_slide(prs, "04. 프로모션 페이지 XPath",
    "TV 세일/프로모션 페이지",
    "bby_tv_pmt1.py",
    ["config_key", "설명", "우선순위", "XPath"],
    pmt_rows, 6, total_pages)

# 7. 트렌딩 딜 XPath
trend_rows = []
for key, val, pri in by_file.get('bby_tv_trend_crawl', []):
    desc = {
        'trending_section': '트렌딩 섹션 영역',
        'tvs_button': 'TV 카테고리 버튼',
        'product_items': '상품 목록',
        'product_item': '개별 상품',
        'product_name': '상품명',
        'product_url': '상품 URL',
        'rank': '트렌딩 순위',
    }.get(key, key)
    val_display = val[:55] + "..." if len(val) > 55 else val
    trend_rows.append((key, desc, f"p{pri}", val_display))

add_xpath_table_slide(prs, "05. 트렌딩 딜 XPath",
    "홈페이지 Trending Deals 섹션",
    "bby_tv_trend_crawl.py",
    ["config_key", "설명", "우선순위", "XPath"],
    trend_rows, 7, total_pages)

# 8. 상세 페이지 XPath
dt_rows = []
for key, val, pri in by_file.get('bby_tv_dt1', []):
    desc = {
        'product_title': '상품명',
        'final_price': '최종 가격',
        'model_number': '모델 번호',
        'overall_rating': '평점',
        'review_count': '리뷰 수',
        'specs_button': '스펙 버튼',
        'see_all_reviews_btn': '리뷰 더보기',
        'next_page_btn': '다음 페이지',
        'top_mentions': '리뷰 키워드',
        'compare_product_title': '비교 상품명',
    }.get(key, key)
    val_display = val[:55] + "..." if len(val) > 55 else val
    dt_rows.append((key, desc, f"p{pri}", val_display))

add_xpath_table_slide(prs, "06. 상세 페이지 XPath",
    "개별 상품 상세 정보",
    "bby_tv_dt1.py",
    ["config_key", "설명", "우선순위", "XPath"],
    dt_rows, 8, total_pages)

# 9. 자주 변경되는 XPath
add_frequency_slide(prs, 9, total_pages)

# 10. 수정 예시
add_example_slide(prs, 10, total_pages)

# 11. Quick Reference
add_quick_ref_slide(prs, 11, total_pages)

# 저장
output_path = 'BestBuy_TV_XPath_Guide_v2.pptx'
prs.save(output_path)
print(f'[OK] 고급 버전 PPT 생성 완료: {output_path}')
